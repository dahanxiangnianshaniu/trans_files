"""
红线条目 11.1.1-1+2 隐私声明与个人数据说明合规检查 — MCP Server

文档发现、读取、检索、表格提取、章节提取。
纯基础设施层，不含任何 LLM 调用。智能分析由 Agent 自身完成。

支持格式：.md .txt .csv .html .xml .json .yaml .yml .rst .ini .cfg .conf
           .pdf .docx .xlsx .xls .pptx .chm（需7-Zip）
"""

import os
import re
import json
import shutil
import tempfile
import subprocess
from typing import Optional

from fastmcp import FastMCP

mcp = FastMCP("redline-doc-audit")

SUPPORTED_TEXT_EXT = {
    '.md', '.txt', '.csv', '.html', '.htm', '.xml',
    '.json', '.yaml', '.yml', '.rst', '.ini', '.cfg', '.conf',
}
SUPPORTED_DOC_EXT = {'.pdf', '.docx', '.doc', '.xlsx', '.xls', '.pptx', '.chm'}
SUPPORTED_ALL_EXT = SUPPORTED_TEXT_EXT | SUPPORTED_DOC_EXT

# ─────────────────────────── CHM 处理 ───────────────────────────

def _extract_chm(chm_path: str, seven_zip_path: str) -> Optional[str]:
    """用7-Zip解压CHM到临时目录，返回临时目录路径。失败返回None。"""
    if not seven_zip_path or not os.path.isfile(seven_zip_path):
        return None
    temp_dir = tempfile.mkdtemp(prefix="chm_extract_")
    try:
        cmd = [seven_zip_path, "x", f"-o{temp_dir}", "-y", chm_path]
        result = subprocess.run(cmd, capture_output=True, timeout=30)
        if result.returncode != 0:
            shutil.rmtree(temp_dir, ignore_errors=True)
            return None
        return temp_dir
    except Exception:
        shutil.rmtree(temp_dir, ignore_errors=True)
        return None


def _read_html_from_dir(temp_dir: str, max_length: int = 32768, start_offset: int = 0) -> str:
    """从解压后的CHM目录中读取所有HTML文件内容，拼接为文本。"""
    texts = []
    total_len = 0
    for root, dirs, files in os.walk(temp_dir):
        for fname in files:
            if fname.lower().endswith(('.html', '.htm')):
                fpath = os.path.join(root, fname)
                try:
                    with open(fpath, 'r', encoding='utf-8', errors='replace') as f:
                        raw = f.read()
                    text = _strip_html_tags(raw)
                    texts.append(text)
                    total_len += len(text)
                    if total_len > max_length * 2:
                        break
                except Exception:
                    continue
        if total_len > max_length * 2:
            break
    full = "\n\n".join(texts)
    if start_offset > 0:
        full = full[start_offset:]
    if len(full) > max_length:
        full = full[:max_length]
    return full


def _strip_html_tags(html: str) -> str:
    """简单去除HTML标签，保留文本内容。"""
    text = re.sub(r'<script[^>]*>.*?</script>', '', html, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r'<br\s*/?\s*>', '\n', text, flags=re.IGNORECASE)
    text = re.sub(r'<p[^>]*>', '\n', text, flags=re.IGNORECASE)
    text = re.sub(r'<li[^>]*>', '\n- ', text, flags=re.IGNORECASE)
    text = re.sub(r'<[^>]+>', '', text)
    text = re.sub(r'&nbsp;', ' ', text)
    text = re.sub(r'&lt;', '<', text)
    text = re.sub(r'&gt;', '>', text)
    text = re.sub(r'&amp;', '&', text)
    text = re.sub(r'&#\d+;', '', text)
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text.strip()


def _cleanup_chm_temp(temp_dir: str):
    if temp_dir and os.path.isdir(temp_dir):
        shutil.rmtree(temp_dir, ignore_errors=True)


# ─────────────────────────── 文档读取 ───────────────────────────

def _read_text_file(path: str, max_length: int = 32768, start_offset: int = 0) -> dict:
    try:
        with open(path, 'r', encoding='utf-8', errors='replace') as f:
            content = f.read()
        total_length = len(content)
        if start_offset > 0:
            content = content[start_offset:]
        truncated = len(content) > max_length
        if truncated:
            content = content[:max_length]
        return {"content": content, "total_length": total_length,
                "read_length": len(content), "truncated": truncated}
    except Exception as e:
        return {"content": "", "total_length": 0, "read_length": 0,
                "truncated": False, "error": str(e)}


def _read_docx(path: str, max_length: int = 32768, start_offset: int = 0) -> dict:
    try:
        from docx import Document
        doc = Document(path)
        parts = []
        for para in doc.paragraphs:
            parts.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                parts.append(row_text)
        content = "\n".join(parts)
        total_length = len(content)
        has_tables = len(doc.tables) > 0
        sections = [p.text for p in doc.paragraphs if p.style and p.style.name.startswith('Heading')]
        if start_offset > 0:
            content = content[start_offset:]
        truncated = len(content) > max_length
        if truncated:
            content = content[:max_length]
        return {"content": content, "total_length": total_length,
                "read_length": len(content), "truncated": truncated,
                "metadata": {"has_tables": has_tables, "sections": sections[:20]}}
    except Exception as e:
        return {"content": "", "total_length": 0, "read_length": 0,
                "truncated": False, "error": str(e)}


def _read_xlsx(path: str, max_length: int = 32768, start_offset: int = 0) -> dict:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        parts = []
        sheets = wb.sheetnames
        for sname in sheets:
            ws = wb[sname]
            parts.append(f"=== Sheet: {sname} ===")
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(str(c) if c is not None else '' for c in row)
                parts.append(row_text)
        wb.close()
        content = "\n".join(parts)
        total_length = len(content)
        if start_offset > 0:
            content = content[start_offset:]
        truncated = len(content) > max_length
        if truncated:
            content = content[:max_length]
        return {"content": content, "total_length": total_length,
                "read_length": len(content), "truncated": truncated,
                "metadata": {"sheets": sheets}}
    except Exception as e:
        return {"content": "", "total_length": 0, "read_length": 0,
                "truncated": False, "error": str(e)}


def _read_pdf(path: str, max_length: int = 32768, start_offset: int = 0) -> dict:
    try:
        import pdfplumber
        parts = []
        page_count = 0
        has_tables = False
        with pdfplumber.open(path) as pdf:
            page_count = len(pdf.pages)
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    parts.append(text)
                tables = page.extract_tables()
                if tables:
                    has_tables = True
                    for table in tables:
                        for row in table:
                            row_text = " | ".join(str(c) if c else '' for c in row)
                            parts.append(row_text)
        content = "\n\n".join(parts)
        total_length = len(content)
        if start_offset > 0:
            content = content[start_offset:]
        truncated = len(content) > max_length
        if truncated:
            content = content[:max_length]
        return {"content": content, "total_length": total_length,
                "read_length": len(content), "truncated": truncated,
                "metadata": {"page_count": page_count, "has_tables": has_tables}}
    except Exception as e:
        return {"content": "", "total_length": 0, "read_length": 0,
                "truncated": False, "error": str(e)}


def _read_pptx(path: str, max_length: int = 32768, start_offset: int = 0) -> dict:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        parts = []
        page_count = len(prs.slides)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text)
        content = "\n\n".join(parts)
        total_length = len(content)
        if start_offset > 0:
            content = content[start_offset:]
        truncated = len(content) > max_length
        if truncated:
            content = content[:max_length]
        return {"content": content, "total_length": total_length,
                "read_length": len(content), "truncated": truncated,
                "metadata": {"page_count": page_count}}
    except Exception as e:
        return {"content": "", "total_length": 0, "read_length": 0,
                "truncated": False, "error": str(e)}


def _read_chm(path: str, seven_zip_path: str, max_length: int = 32768, start_offset: int = 0) -> dict:
    temp_dir = _extract_chm(path, seven_zip_path)
    if not temp_dir:
        return {"content": "", "total_length": 0, "read_length": 0,
                "truncated": False, "error": "CHM解压失败，请检查7-Zip路径"}
    try:
        content = _read_html_from_dir(temp_dir, max_length * 2, start_offset)
        total_length = len(content)
        truncated = len(content) > max_length
        if truncated:
            content = content[:max_length]
        return {"content": content, "total_length": total_length,
                "read_length": len(content), "truncated": truncated}
    finally:
        _cleanup_chm_temp(temp_dir)


def _read_any(path: str, max_length: int = 32768, start_offset: int = 0,
              seven_zip_path: str = "") -> dict:
    """统一文档读取入口，根据扩展名路由到对应reader。"""
    ext = os.path.splitext(path)[1].lower()
    result = {"file_path": path, "format": ext}

    if ext in SUPPORTED_TEXT_EXT:
        r = _read_text_file(path, max_length, start_offset)
    elif ext == '.docx':
        r = _read_docx(path, max_length, start_offset)
    elif ext in ('.xlsx', '.xls'):
        r = _read_xlsx(path, max_length, start_offset)
    elif ext == '.pdf':
        r = _read_pdf(path, max_length, start_offset)
    elif ext == '.pptx':
        r = _read_pptx(path, max_length, start_offset)
    elif ext == '.chm':
        r = _read_chm(path, seven_zip_path, max_length, start_offset)
    elif ext == '.doc':
        result.update({"content": "", "total_length": 0, "read_length": 0,
                        "truncated": False, "error": ".doc格式需先转换为.docx"})
        return result
    else:
        result.update({"content": "", "total_length": 0, "read_length": 0,
                        "truncated": False, "error": f"不支持的格式: {ext}"})
        return result

    result.update(r)
    if "format" not in result:
        result["format"] = ext
    return result


# ─────────────────────────── 表格提取 ───────────────────────────

def _extract_tables_docx(path: str, table_index: Optional[int] = None) -> dict:
    try:
        from docx import Document
        doc = Document(path)
        tables_data = []
        for i, table in enumerate(doc.tables):
            if table_index is not None and i != table_index:
                continue
            headers = [cell.text.strip() for cell in table.rows[0].cells] if table.rows else []
            rows = []
            empty_cells = []
            for ri, row in enumerate(table.rows[1:], start=2):
                row_dict = {}
                for ci, cell in enumerate(row.cells):
                    val = cell.text.strip()
                    col_name = headers[ci] if ci < len(headers) else f"col_{ci}"
                    row_dict[col_name] = val
                    if not val:
                        empty_cells.append({"row": ri, "column": col_name, "value": ""})
                rows.append(row_dict)
            tables_data.append({
                "table_index": i, "sheet_name": None,
                "headers": headers, "row_count": len(rows),
                "rows": rows, "empty_cells": empty_cells,
            })
        return {"source_format": "docx", "tables": tables_data, "total_tables_found": len(doc.tables)}
    except Exception as e:
        return {"source_format": "docx", "tables": [], "total_tables_found": 0, "error": str(e)}


def _extract_tables_xlsx(path: str, table_index: Optional[int] = None,
                          sheet_name: Optional[str] = None, max_rows: int = 500) -> dict:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        sheets = wb.sheetnames
        tables_data = []
        target_sheets = [sheet_name] if sheet_name else sheets
        for sname in target_sheets:
            if sname not in sheets:
                continue
            ws = wb[sname]
            rows_list = list(ws.iter_rows(values_only=True))
            if not rows_list:
                continue
            headers = [str(c) if c is not None else '' for c in rows_list[0]]
            data_rows = rows_list[1:max_rows + 1]
            rows = []
            empty_cells = []
            for ri, row in enumerate(data_rows, start=2):
                row_dict = {}
                for ci, cell in enumerate(row):
                    val = str(cell) if cell is not None else ''
                    col_name = headers[ci] if ci < len(headers) else f"col_{ci}"
                    row_dict[col_name] = val
                    if not val:
                        empty_cells.append({"row": ri, "column": col_name, "value": ""})
                rows.append(row_dict)
            tables_data.append({
                "table_index": 0, "sheet_name": sname,
                "headers": headers, "row_count": len(rows),
                "rows": rows, "empty_cells": empty_cells,
            })
        wb.close()
        total = len(sheets)
        return {"source_format": "xlsx", "tables": tables_data, "total_tables_found": total}
    except Exception as e:
        return {"source_format": "xlsx", "tables": [], "total_tables_found": 0, "error": str(e)}


def _extract_tables_pdf(path: str, table_index: Optional[int] = None,
                         max_rows: int = 500) -> dict:
    try:
        import pdfplumber
        tables_data = []
        with pdfplumber.open(path) as pdf:
            ti = 0
            for page in pdf.pages:
                tables = page.extract_tables()
                for table in tables:
                    if table_index is not None and ti != table_index:
                        ti += 1
                        continue
                    if not table:
                        ti += 1
                        continue
                    headers = [str(c) if c else '' for c in table[0]]
                    data_rows = table[1:max_rows + 1]
                    rows = []
                    empty_cells = []
                    for ri, row in enumerate(data_rows, start=2):
                        row_dict = {}
                        for ci, cell in enumerate(row):
                            val = str(cell) if cell else ''
                            col_name = headers[ci] if ci < len(headers) else f"col_{ci}"
                            row_dict[col_name] = val
                            if not val:
                                empty_cells.append({"row": ri, "column": col_name, "value": ""})
                        rows.append(row_dict)
                    tables_data.append({
                        "table_index": ti, "sheet_name": None,
                        "headers": headers, "row_count": len(rows),
                        "rows": rows, "empty_cells": empty_cells,
                    })
                    ti += 1
        return {"source_format": "pdf", "tables": tables_data, "total_tables_found": ti}
    except Exception as e:
        return {"source_format": "pdf", "tables": [], "total_tables_found": 0, "error": str(e)}


def _parse_md_tables(content: str, table_index: Optional[int] = None,
                      max_rows: int = 500) -> list:
    lines = content.split('\n')
    tables = []
    i = 0
    ti = 0
    while i < len(lines):
        line = lines[i]
        if '|' in line and i + 1 < len(lines) and re.match(r'^[\s|:-]+$', lines[i + 1]):
            if table_index is not None and ti != table_index:
                ti += 1
                i += 2
                continue
            headers = [c.strip() for c in line.split('|')[1:-1]]
            i += 2
            rows = []
            empty_cells = []
            ri = 2
            while i < len(lines) and '|' in lines[i] and not re.match(r'^[\s|:-]+$', lines[i]):
                cells = [c.strip() for c in lines[i].split('|')[1:-1]]
                row_dict = {}
                for ci, val in enumerate(cells):
                    col_name = headers[ci] if ci < len(headers) else f"col_{ci}"
                    row_dict[col_name] = val
                    if not val:
                        empty_cells.append({"row": ri, "column": col_name, "value": ""})
                rows.append(row_dict)
                ri += 1
                i += 1
                if len(rows) >= max_rows:
                    break
            tables.append({
                "table_index": ti, "sheet_name": None,
                "headers": headers, "row_count": len(rows),
                "rows": rows, "empty_cells": empty_cells,
            })
            ti += 1
        else:
            i += 1
    return tables, ti


def _extract_tables_md(path: str, table_index: Optional[int] = None,
                        max_rows: int = 500) -> dict:
    try:
        with open(path, 'r', encoding='utf-8', errors='replace') as f:
            content = f.read()
        tables, total = _parse_md_tables(content, table_index, max_rows)
        return {"source_format": "md", "tables": tables, "total_tables_found": total}
    except Exception as e:
        return {"source_format": "md", "tables": [], "total_tables_found": 0, "error": str(e)}


def _extract_tables_csv(path: str, max_rows: int = 500) -> dict:
    try:
        import csv
        with open(path, 'r', encoding='utf-8', errors='replace', newline='') as f:
            reader = csv.reader(f)
            rows_raw = list(reader)
        if not rows_raw:
            return {"source_format": "csv", "tables": [], "total_tables_found": 0}
        headers = rows_raw[0]
        data = rows_raw[1:max_rows + 1]
        rows = []
        empty_cells = []
        for ri, row in enumerate(data, start=2):
            row_dict = {}
            for ci, val in enumerate(row):
                col_name = headers[ci] if ci < len(headers) else f"col_{ci}"
                row_dict[col_name] = val.strip()
                if not val.strip():
                    empty_cells.append({"row": ri, "column": col_name, "value": ""})
            rows.append(row_dict)
        tables = [{"table_index": 0, "sheet_name": None,
                    "headers": headers, "row_count": len(rows),
                    "rows": rows, "empty_cells": empty_cells}]
        return {"source_format": "csv", "tables": tables, "total_tables_found": 1}
    except Exception as e:
        return {"source_format": "csv", "tables": [], "total_tables_found": 0, "error": str(e)}


def _extract_tables_html(path: str, table_index: Optional[int] = None,
                          max_rows: int = 500) -> dict:
    try:
        from bs4 import BeautifulSoup
        with open(path, 'r', encoding='utf-8', errors='replace') as f:
            content = f.read()
        soup = BeautifulSoup(content, 'html.parser')
        tables_data = []
        for ti, table in enumerate(soup.find_all('table')):
            if table_index is not None and ti != table_index:
                continue
            headers = []
            header_row = table.find('tr')
            if header_row:
                for th in header_row.find_all(['th', 'td']):
                    headers.append(th.get_text(strip=True))
            rows = []
            empty_cells = []
            ri = 2
            for tr in table.find_all('tr')[1:]:
                cells = tr.find_all(['td', 'th'])
                if not cells:
                    continue
                row_dict = {}
                for ci, cell in enumerate(cells):
                    val = cell.get_text(strip=True)
                    col_name = headers[ci] if ci < len(headers) else f"col_{ci}"
                    row_dict[col_name] = val
                    if not val:
                        empty_cells.append({"row": ri, "column": col_name, "value": ""})
                rows.append(row_dict)
                ri += 1
                if len(rows) >= max_rows:
                    break
            tables_data.append({
                "table_index": ti, "sheet_name": None,
                "headers": headers if headers else [],
                "row_count": len(rows), "rows": rows, "empty_cells": empty_cells,
            })
        total = len(soup.find_all('table'))
        return {"source_format": "html", "tables": tables_data, "total_tables_found": total}
    except Exception as e:
        return {"source_format": "html", "tables": [], "total_tables_found": 0, "error": str(e)}


def _extract_tables_chm(path: str, seven_zip_path: str,
                          table_index: Optional[int] = None,
                          max_rows: int = 500) -> dict:
    temp_dir = _extract_chm(path, seven_zip_path)
    if not temp_dir:
        return {"source_format": "chm", "tables": [], "total_tables_found": 0,
                "error": "CHM解压失败"}
    try:
        from bs4 import BeautifulSoup
        all_tables = []
        ti = 0
        for root, dirs, files in os.walk(temp_dir):
            for fname in files:
                if fname.lower().endswith(('.html', '.htm')):
                    fpath = os.path.join(root, fname)
                    try:
                        with open(fpath, 'r', encoding='utf-8', errors='replace') as f:
                            html = f.read()
                        soup = BeautifulSoup(html, 'html.parser')
                        for table in soup.find_all('table'):
                            if table_index is not None and ti != table_index:
                                ti += 1
                                continue
                            headers = []
                            header_row = table.find('tr')
                            if header_row:
                                for th in header_row.find_all(['th', 'td']):
                                    headers.append(th.get_text(strip=True))
                            rows = []
                            empty_cells = []
                            ri = 2
                            for tr in table.find_all('tr')[1:]:
                                cells = tr.find_all(['td', 'th'])
                                if not cells:
                                    continue
                                row_dict = {}
                                for ci, cell in enumerate(cells):
                                    val = cell.get_text(strip=True)
                                    col_name = headers[ci] if ci < len(headers) else f"col_{ci}"
                                    row_dict[col_name] = val
                                    if not val:
                                        empty_cells.append({"row": ri, "column": col_name, "value": ""})
                                rows.append(row_dict)
                                ri += 1
                                if len(rows) >= max_rows:
                                    break
                            all_tables.append({
                                "table_index": ti, "sheet_name": None,
                                "headers": headers, "row_count": len(rows),
                                "rows": rows, "empty_cells": empty_cells,
                            })
                            ti += 1
                    except Exception:
                        continue
        return {"source_format": "chm", "tables": all_tables, "total_tables_found": ti}
    finally:
        _cleanup_chm_temp(temp_dir)


# ─────────────────────────── 章节提取 ───────────────────────────

def _list_sections_docx(path: str) -> dict:
    try:
        from docx import Document
        doc = Document(path)
        sections = []
        idx = 0
        for para in doc.paragraphs:
            if para.style and para.style.name.startswith('Heading'):
                level = int(para.style.name.replace('Heading ', '').replace('Heading', '1') or '1')
                sections.append({"level": level, "title": para.text.strip(), "index": idx})
                idx += 1
        return {"source_format": "docx", "sections": sections, "sheets": None}
    except Exception as e:
        return {"source_format": "docx", "sections": [], "sheets": None, "error": str(e)}


def _list_sections_xlsx(path: str) -> dict:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True)
        sheets = wb.sheetnames
        wb.close()
        return {"source_format": "xlsx", "sections": None, "sheets": sheets}
    except Exception as e:
        return {"source_format": "xlsx", "sections": None, "sheets": [], "error": str(e)}


def _list_sections_pdf(path: str) -> dict:
    try:
        import pdfplumber
        sections = []
        with pdfplumber.open(path) as pdf:
            if pdf.outline:
                idx = 0
                for item in pdf.outline:
                    if isinstance(item, dict) and 'title' in item:
                        sections.append({"level": 1, "title": item['title'], "index": idx})
                        idx += 1
        return {"source_format": "pdf", "sections": sections, "sheets": None}
    except Exception as e:
        return {"source_format": "pdf", "sections": [], "sheets": None, "error": str(e)}


def _list_sections_md(path: str) -> dict:
    try:
        with open(path, 'r', encoding='utf-8', errors='replace') as f:
            content = f.read()
        sections = []
        idx = 0
        for line in content.split('\n'):
            m = re.match(r'^(#{1,6})\s+(.+)$', line)
            if m:
                level = len(m.group(1))
                sections.append({"level": level, "title": m.group(2).strip(), "index": idx})
                idx += 1
        return {"source_format": "md", "sections": sections, "sheets": None}
    except Exception as e:
        return {"source_format": "md", "sections": [], "sheets": None, "error": str(e)}


def _list_sections_html(path: str) -> dict:
    try:
        from bs4 import BeautifulSoup
        with open(path, 'r', encoding='utf-8', errors='replace') as f:
            content = f.read()
        soup = BeautifulSoup(content, 'html.parser')
        sections = []
        idx = 0
        for tag in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
            for elem in soup.find_all(tag):
                level = int(tag[1])
                sections.append({"level": level, "title": elem.get_text(strip=True), "index": idx})
                idx += 1
        return {"source_format": "html", "sections": sections, "sheets": None}
    except Exception as e:
        return {"source_format": "html", "sections": [], "sheets": None, "error": str(e)}


def _list_sections_chm(path: str, seven_zip_path: str) -> dict:
    temp_dir = _extract_chm(path, seven_zip_path)
    if not temp_dir:
        return {"source_format": "chm", "sections": [], "sheets": None, "error": "CHM解压失败"}
    try:
        from bs4 import BeautifulSoup
        all_sections = []
        idx = 0
        for root, dirs, files in os.walk(temp_dir):
            for fname in files:
                if fname.lower().endswith(('.html', '.htm')):
                    fpath = os.path.join(root, fname)
                    try:
                        with open(fpath, 'r', encoding='utf-8', errors='replace') as f:
                            html = f.read()
                        soup = BeautifulSoup(html, 'html.parser')
                        for tag in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                            for elem in soup.find_all(tag):
                                level = int(tag[1])
                                all_sections.append({"level": level, "title": elem.get_text(strip=True), "index": idx})
                                idx += 1
                    except Exception:
                        continue
        return {"source_format": "chm", "sections": all_sections, "sheets": None}
    finally:
        _cleanup_chm_temp(temp_dir)


# ─────────────────────────── 文档评分 ───────────────────────────

def _score_document(filename: str, content_preview: str, keywords: list) -> int:
    """3信号评分：filename 50% + content_preview 30% + structural 20%"""
    fn_lower = filename.lower()

    # filename_score
    fn_score = 0
    for kw in keywords:
        kw_lower = kw.lower()
        if kw_lower in fn_lower:
            fn_score = max(fn_score, 100)
            break
    if fn_score == 0:
        for kw in keywords:
            kw_lower = kw.lower()
            if any(part in fn_lower for part in kw_lower.split() if len(part) > 1):
                fn_score = max(fn_score, 60)
    if fn_score == 0:
        related_words = ['声明', '政策', '隐私', 'privacy', 'data', 'personal', '说明', '清单']
        for w in related_words:
            if w.lower() in fn_lower:
                fn_score = max(fn_score, 30)
                break

    # content_preview_score
    cp_score = 0
    if content_preview:
        preview_lower = content_preview.lower()
        hits = sum(1 for kw in keywords if kw.lower() in preview_lower)
        if hits > 0:
            cp_score = min(100, int(hits / len(keywords) * 100) + 20)

    # structural_score
    st_score = 0
    if content_preview:
        preview_lower = content_preview.lower()
        structural_keywords = ['个人数据', '数据类型', '收集目的', '处理方式', '存留期限',
                               'personal data', 'data type', 'purpose', 'retention']
        for sk in structural_keywords:
            if sk.lower() in preview_lower:
                st_score = max(st_score, 80)
                break

    total = int(fn_score * 0.5 + cp_score * 0.3 + st_score * 0.2)
    return total


# ═══════════════════════════ MCP 工具 ═══════════════════════════

@mcp.tool()
def scan_documents(
    doc_path: str,
    keywords: list,
    search_mode: str = "precise",
    file_types: list = None,
    seven_zip_path: str = "",
) -> dict:
    """扫描目录发现相关文档。支持precise/broad/fullscan三种搜索模式。

    Args:
        doc_path: 文档根目录
        keywords: 搜索关键词列表
        search_mode: 搜索模式 - precise(精确,默认)/broad(扩展)/fullscan(全量扫描)
        file_types: 文件扩展名过滤(如['.pdf','.docx'])，为空则包含所有支持的格式
        seven_zip_path: 7-Zip可执行文件路径，用于.chm文件
    """
    if not os.path.isdir(doc_path):
        return {"success": False, "error": f"目录不存在: {doc_path}"}

    all_files = []
    for root, dirs, files in os.walk(doc_path):
        for fname in files:
            ext = os.path.splitext(fname)[1].lower()
            if file_types:
                if ext not in [ft.lower() for ft in file_types]:
                    continue
            else:
                if ext not in SUPPORTED_ALL_EXT:
                    continue
            fpath = os.path.join(root, fname)
            try:
                fsize = os.path.getsize(fpath)
            except OSError:
                fsize = 0
            all_files.append({"file_path": fpath, "file_name": fname,
                              "file_type": ext, "file_size": fsize})

    if search_mode == "fullscan":
        matched = []
        for f in all_files:
            preview = _get_content_preview(f["file_path"], f["file_type"], 200, seven_zip_path)
            matched.append({
                "file_path": f["file_path"], "file_name": f["file_name"],
                "file_type": f["file_type"], "file_size": f["file_size"],
                "match_score": 0, "match_signals": {},
                "content_preview": preview,
            })
        return {
            "success": True,
            "total_files_scanned": len(all_files),
            "search_mode": search_mode,
            "matched_files": matched,
            "unmatched_file_list": [],
        }

    # precise / broad 模式
    threshold = 50 if search_mode == "precise" else 30
    matched = []
    unmatched = []
    for f in all_files:
        preview = _get_content_preview(f["file_path"], f["file_type"], 500, seven_zip_path)
        score = _score_document(f["file_name"], preview, keywords)

        signals = {"filename_match": [], "content_preview_match": []}
        fn_lower = f["file_name"].lower()
        for kw in keywords:
            if kw.lower() in fn_lower:
                signals["filename_match"].append(kw)
            if preview and kw.lower() in preview.lower():
                signals["content_preview_match"].append(kw)

        if score >= threshold:
            matched.append({
                "file_path": f["file_path"], "file_name": f["file_name"],
                "file_type": f["file_type"], "file_size": f["file_size"],
                "match_score": score, "match_signals": signals,
                "content_preview": preview[:500] if preview else "",
            })
        else:
            unmatched.append(f["file_name"])

    matched.sort(key=lambda x: x["match_score"], reverse=True)
    return {
        "success": True,
        "total_files_scanned": len(all_files),
        "search_mode": search_mode,
        "matched_files": matched,
        "unmatched_file_list": unmatched[:50],
    }


def _get_content_preview(path: str, ext: str, chars: int = 500,
                          seven_zip_path: str = "") -> str:
    """获取文件内容预览（前N字符）。"""
    try:
        if ext in SUPPORTED_TEXT_EXT:
            with open(path, 'r', encoding='utf-8', errors='replace') as f:
                return f.read(chars)
        elif ext == '.docx':
            from docx import Document
            doc = Document(path)
            return " ".join(p.text for p in doc.paragraphs[:10])[:chars]
        elif ext in ('.xlsx', '.xls'):
            import openpyxl
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            parts = [f"页签: {', '.join(wb.sheetnames[:5])}"]
            if wb.sheetnames:
                ws = wb[wb.sheetnames[0]]
                for row in ws.iter_rows(values_only=True, max_row=5):
                    parts.append(" | ".join(str(c) if c else '' for c in row))
            wb.close()
            return " ".join(parts)[:chars]
        elif ext == '.pdf':
            import pdfplumber
            with pdfplumber.open(path) as pdf:
                if pdf.pages:
                    return (pdf.pages[0].extract_text() or "")[:chars]
        elif ext == '.pptx':
            from pptx import Presentation
            prs = Presentation(path)
            parts = []
            for slide in prs.slides[:3]:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        parts.append(shape.text)
            return " ".join(parts)[:chars]
        elif ext == '.chm':
            temp_dir = _extract_chm(path, seven_zip_path)
            if temp_dir:
                try:
                    return _read_html_from_dir(temp_dir, chars)[:chars]
                finally:
                    _cleanup_chm_temp(temp_dir)
    except Exception:
        pass
    return ""


@mcp.tool()
def read_document(
    file_path: str,
    max_length: int = 32768,
    start_offset: int = 0,
    seven_zip_path: str = "",
) -> dict:
    """读取文档全文内容，自动适配10+格式(.pdf/.docx/.xlsx/.chm/.html/.md等)。

    Args:
        file_path: 文件路径
        max_length: 最大字符数，默认32768
        start_offset: 起始偏移字符，默认0
        seven_zip_path: 7-Zip路径(.chm文件需要)
    """
    if not os.path.isfile(file_path):
        return {"success": False, "error": f"文件不存在: {file_path}"}

    result = _read_any(file_path, max_length, start_offset, seven_zip_path)
    result["success"] = "error" not in result
    return result


@mcp.tool()
def search_in_document(
    file_path: str,
    keywords: list,
    context_lines: int = 5,
    max_results: int = 50,
    seven_zip_path: str = "",
) -> dict:
    """文档内关键词检索，返回行号+上下文。替代RAG语义检索。

    Args:
        file_path: 文件路径
        keywords: 搜索关键词列表，支持多关键词并行搜索
        context_lines: 上下文行数，默认5
        max_results: 最大返回数，默认50
        seven_zip_path: 7-Zip路径(.chm文件需要)
    """
    if not os.path.isfile(file_path):
        return {"success": False, "error": f"文件不存在: {file_path}"}

    # 读取全文
    read_result = _read_any(file_path, 500000, 0, seven_zip_path)
    content = read_result.get("content", "")
    if not content:
        return {"success": False, "error": read_result.get("error", "无法读取文件内容")}

    lines = content.split('\n')
    results = []
    total_matches = 0

    for kw in keywords:
        kw_lower = kw.lower()
        count = 0
        for li, line in enumerate(lines):
            if kw_lower in line.lower():
                total_matches += 1
                count += 1
                if count > max_results:
                    continue
                before = [f"L{li - context_lines + j}: {lines[li - context_lines + j]}"
                          for j in range(context_lines) if li - context_lines + j >= 0]
                after = [f"L{li + 1 + j}: {lines[li + 1 + j]}"
                         for j in range(context_lines) if li + 1 + j < len(lines)]
                results.append({
                    "keyword": kw,
                    "line_number": li + 1,
                    "context_before": before,
                    "matched_line": line.strip(),
                    "context_after": after,
                })

    return {"success": True, "total_matches": total_matches, "results": results[:max_results]}


@mcp.tool()
def extract_tables(
    file_path: str,
    table_index: int = -1,
    sheet_name: str = "",
    max_rows: int = 500,
    seven_zip_path: str = "",
) -> dict:
    """通用表格结构提取，不限Excel，Word/PDF/HTML/CHM/Markdown中的表格均可提取。
    统一输出headers+rows+empty_cells结构。

    Args:
        file_path: 文件路径
        table_index: 指定第几个表格(0起)，-1表示提取全部
        sheet_name: Excel页签名，仅.xlsx/.xls有效
        max_rows: 每个表格最大行数，默认500
        seven_zip_path: 7-Zip路径(.chm文件需要)
    """
    if not os.path.isfile(file_path):
        return {"success": False, "error": f"文件不存在: {file_path}"}

    ext = os.path.splitext(file_path)[1].lower()
    ti = table_index if table_index >= 0 else None

    if ext == '.docx':
        result = _extract_tables_docx(file_path, ti)
    elif ext in ('.xlsx', '.xls'):
        sn = sheet_name if sheet_name else None
        result = _extract_tables_xlsx(file_path, ti, sn, max_rows)
    elif ext == '.pdf':
        result = _extract_tables_pdf(file_path, ti, max_rows)
    elif ext == '.chm':
        result = _extract_tables_chm(file_path, seven_zip_path, ti, max_rows)
    elif ext in ('.html', '.htm'):
        result = _extract_tables_html(file_path, ti, max_rows)
    elif ext == '.md':
        result = _extract_tables_md(file_path, ti, max_rows)
    elif ext == '.csv':
        result = _extract_tables_csv(file_path, max_rows)
    else:
        return {"success": False, "error": f"不支持从 {ext} 格式提取表格"}

    result["success"] = "error" not in result
    return result


@mcp.tool()
def list_sections(
    file_path: str,
    seven_zip_path: str = "",
) -> dict:
    """文档章节/目录/页签结构提取，帮助快速定位相关内容。

    Args:
        file_path: 文件路径
        seven_zip_path: 7-Zip路径(.chm文件需要)
    """
    if not os.path.isfile(file_path):
        return {"success": False, "error": f"文件不存在: {file_path}"}

    ext = os.path.splitext(file_path)[1].lower()

    if ext == '.docx':
        result = _list_sections_docx(file_path)
    elif ext in ('.xlsx', '.xls'):
        result = _list_sections_xlsx(file_path)
    elif ext == '.pdf':
        result = _list_sections_pdf(file_path)
    elif ext == '.chm':
        result = _list_sections_chm(file_path, seven_zip_path)
    elif ext == '.md':
        result = _list_sections_md(file_path)
    elif ext in ('.html', '.htm'):
        result = _list_sections_html(file_path)
    elif ext in ('.txt', '.rst'):
        result = {"source_format": ext, "sections": [], "sheets": None}
    else:
        return {"success": False, "error": f"不支持从 {ext} 格式提取章节"}

    result["success"] = "error" not in result
    return result


# ═══════════════════════════ 入口 ═══════════════════════════

if __name__ == "__main__":
    mcp.run()
